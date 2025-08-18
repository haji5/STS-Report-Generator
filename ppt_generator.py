from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import glob
import os
from charts import generate_grouped_barplot
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image


# Constants for slide dimensions and positioning
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
CONTENT_LEFT = Inches(2.0)
CONTENT_TOP = Inches(0.5)
CONTENT_WIDTH = Inches(9.33)
TITLE_FONT_SIZE = Pt(36)
SUBTITLE_FONT_SIZE = Pt(18)
HEADER_FONT_SIZE = Pt(32)
LABEL_FONT_SIZE = Pt(16)

# Chart positioning constants
MAX_CHART_WIDTH = Inches(5.5)
MAX_BARPLOT_WIDTH = Inches(6.5)
PIE_CHART_WIDTH = Inches(4.0)
VERTICAL_SPACING = Inches(0.3)

# File name cleaning characters
FILENAME_REPLACE_CHARS = [' ', ',', '(', ')']


def clean_filename(name):
    """Clean DMO name for use in filenames."""
    safe_name = name
    for char in FILENAME_REPLACE_CHARS:
        safe_name = safe_name.replace(char, '_' if char == ' ' else '')
    return safe_name


def get_dmo_name(slide_info):
    """Extract DMO name from slide info, preferring 'DMO' column if available."""
    return slide_info.get("DMO", slide_info["dmo_name"])


def remove_placeholders(slide):
    """Remove all default placeholders from a slide."""
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            slide.shapes._spTree.remove(shape._element)


def add_textbox_with_formatting(slide, left, top, width, height, text, font_size, bold=False):
    """Add a formatted textbox to a slide."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = font_size
    text_frame.paragraphs[0].font.bold = bold
    return textbox


def add_image_with_aspect_ratio(slide, image_path, left, top, max_width, max_height=None):
    """Add an image to slide while preserving aspect ratio."""
    try:
        with Image.open(image_path) as img:
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height

            desired_width = max_width
            desired_height = desired_width / aspect_ratio

            if max_height and desired_height > max_height:
                desired_height = max_height
                desired_width = desired_height * aspect_ratio

            # Center the image horizontally within the allocated space
            adjusted_left = left + (max_width - desired_width) / 2

            slide.shapes.add_picture(image_path, adjusted_left, top,
                                     width=desired_width, height=desired_height)
            return True
    except Exception as e:
        print(f"[ERROR] Failed to process image {image_path}: {e}")
        # Fallback to adding without aspect ratio preservation
        slide.shapes.add_picture(image_path, left, top, width=max_width)
        return False


def find_chart_files(dmo_name, chart_type):
    """Find chart files for a given DMO and chart type."""
    safe_name = clean_filename(dmo_name)

    # Define potential file patterns
    patterns = {
        'barplot': [
            f"{safe_name}_top5_origins.png",
            f"{safe_name}*_barplot*.png",
            f"{safe_name}*origins*.png"
        ],
        'pie': [
            f"{safe_name}_province_pie.png",
            f"{safe_name}*_pie*.png"
        ],
        'outprov': [
            f"{safe_name}_top5_outofprov_origins.png",
            f"{safe_name}*outprov*.png",
            f"{safe_name}*out-of-prov*.png"
        ]
    }

    for pattern in patterns.get(chart_type, []):
        if '*' in pattern:
            # Use glob for wildcard patterns
            matches = glob.glob(os.path.join('output', pattern))
            if matches:
                return matches[0]
        else:
            # Direct file check
            file_path = os.path.join('output', pattern)
            if os.path.exists(file_path):
                return file_path

    return None


def create_presentation(title, slides, ppt_path):
    """Create a PowerPoint presentation from slide data."""
    print("[DEBUG] Starting create_presentation")
    prs = Presentation()

    # Set widescreen format (16:9)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Get available layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[6]  # Blank content slide

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    # Collect comprehensive slide information for contents
    content_entries = []
    slide_number = 2  # Start after title slide

    # Extract years and DMO names from slides
    year_from_slides = None
    comp_year_from_slides = None
    dmo_names = set()

    for slide_info in slides:
        dmo_name = get_dmo_name(slide_info)
        dmo_names.add(dmo_name)
        year_from_slides = slide_info["year"]
        if slide_info.get("comp_year"):
            comp_year_from_slides = slide_info["comp_year"]

    dmo_names = sorted(list(dmo_names))

    # Add contents slide placeholder - we'll populate it after counting all slides
    contents_slide_index = len(prs.slides)

    # Group slides by DMO and type for better organization
    dmo_slide_groups = {}
    for slide_info in slides:
        dmo_name = get_dmo_name(slide_info)
        slide_type = slide_info.get("slide_type", "visitors_trips")

        if dmo_name not in dmo_slide_groups:
            dmo_slide_groups[dmo_name] = {}
        if slide_type not in dmo_slide_groups[dmo_name]:
            dmo_slide_groups[dmo_name][slide_type] = []
        dmo_slide_groups[dmo_name][slide_type].append(slide_info)

    # Count and organize all slides for contents
    current_page = 3  # After title and contents

    # 1. Individual DMO slides (Visitors & Trips, Nights)
    for dmo_name in sorted(dmo_names):
        if dmo_name in dmo_slide_groups:
            # Visitors & Trips slide
            if "visitors_trips" in dmo_slide_groups[dmo_name]:
                if comp_year_from_slides:
                    content_entries.append({
                        "title": f"{dmo_name} - Visitors & Trips ({comp_year_from_slides} vs {year_from_slides})",
                        "page": current_page,
                        "category": "Individual Destinations"
                    })
                else:
                    content_entries.append({
                        "title": f"{dmo_name} - Visitors & Trips ({year_from_slides})",
                        "page": current_page,
                        "category": "Individual Destinations"
                    })
                current_page += 1

            # Nights slide
            if "nights" in dmo_slide_groups[dmo_name]:
                if comp_year_from_slides:
                    content_entries.append({
                        "title": f"{dmo_name} - Overnight Stays ({comp_year_from_slides} vs {year_from_slides})",
                        "page": current_page,
                        "category": "Individual Destinations"
                    })
                else:
                    content_entries.append({
                        "title": f"{dmo_name} - Overnight Stays ({year_from_slides})",
                        "page": current_page,
                        "category": "Individual Destinations"
                    })
                current_page += 1

    # 2. Top 5 Origins Analysis slides
    for dmo_name in sorted(dmo_names):
        safe_name = clean_filename(dmo_name)

        # Check if origin charts exist for this DMO
        top5_origins_path = os.path.join('output', f"{safe_name}_top5_origins.png")
        if os.path.exists(top5_origins_path):
            content_entries.append({
                "title": f"{dmo_name} - Top 5 Origins by Quarter",
                "page": current_page,
                "category": "Origins Analysis"
            })
            current_page += 1

    # 3. PRIZM & Traveller Segmentation slides
    for dmo_name in sorted(dmo_names):
        safe_name = clean_filename(dmo_name)

        # Check if PRIZM charts exist for this DMO
        circular_path = os.path.join('output', f"{safe_name}_prizm_circular_barplot_{year_from_slides}.png")
        if os.path.exists(circular_path):
            if comp_year_from_slides:
                content_entries.append({
                    "title": f"{dmo_name} - PRIZM & Traveller Segmentation ({comp_year_from_slides} vs {year_from_slides})",
                    "page": current_page,
                    "category": "PRIZM Analysis"
                })
            else:
                content_entries.append({
                    "title": f"{dmo_name} - PRIZM & Traveller Segmentation ({year_from_slides})",
                    "page": current_page,
                    "category": "PRIZM Analysis"
                })
            current_page += 1

    # 4. Quarterly PRIZM Comparison slides (if comparison year exists)
    if comp_year_from_slides:
        for dmo_name in sorted(dmo_names):
            safe_name = clean_filename(dmo_name)
            quarterly_prizm_path = os.path.join('output', f"{safe_name}_quarterly_prizm_barplots_{comp_year_from_slides}_{year_from_slides}.png")

            if os.path.exists(quarterly_prizm_path):
                content_entries.append({
                    "title": f"{dmo_name} - Quarterly PRIZM Comparison ({comp_year_from_slides} vs {year_from_slides})",
                    "page": current_page,
                    "category": "Quarterly Analysis"
                })
                current_page += 1

    # Create contents slides (split into multiple if needed)
    max_entries_per_slide = 16
    total_entries = len(content_entries)
    num_contents_slides = (total_entries + max_entries_per_slide - 1) // max_entries_per_slide

    for contents_page in range(num_contents_slides):
        contents_slide = prs.slides.add_slide(content_slide_layout)
        remove_placeholders(contents_slide)

        # Title
        title_text = "Contents" if num_contents_slides == 1 else f"Contents ({contents_page + 1} of {num_contents_slides})"
        left = Inches(1.0)
        top = Inches(0.3)
        width = Inches(11.33)
        height = Inches(0.8)

        title_box = contents_slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = 1  # Center

        # Contents list
        start_idx = contents_page * max_entries_per_slide
        end_idx = min(start_idx + max_entries_per_slide, total_entries)

        content_top = Inches(1.2)
        content_left = Inches(1.5)
        content_width = Inches(10.33)
        content_height = Inches(5.8)

        textbox = contents_slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        tf = textbox.text_frame
        tf.clear()

        current_category = ""
        for i in range(start_idx, end_idx):
            entry = content_entries[i]

            # Add category header if it's a new category
            if entry["category"] != current_category:
                current_category = entry["category"]

                # Add some spacing before new category (except for first entry)
                if i > start_idx:
                    p = tf.add_paragraph()
                    p.text = ""

                # Category header
                p = tf.add_paragraph()
                p.text = current_category
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 86, 145)  # Blue color

                # Add spacing after category header
                p = tf.add_paragraph()
                p.text = ""

            # Content entry
            p = tf.add_paragraph()
            p.text = f"{entry['title']} {'.' * (80 - len(entry['title']))} {entry['page']}"
            p.font.size = Pt(14)
            p.level = 1

    for slide_info in slides:
        dmo_name = slide_info["dmo_name"]
        # Use DMO column value as the name if available
        if "DMO" in slide_info:
            dmo_name = slide_info["DMO"]
        year = slide_info["year"]
        charts = slide_info["charts"]
        slide_type = slide_info.get("slide_type", "visitors_trips")

        slide = prs.slides.add_slide(content_slide_layout)
        # Remove all default placeholders (like 'Click to add text')
        remove_placeholders(slide)

        # Better positioning for widescreen format
        left = Inches(1.0)  # Increased from 0.5 for better centering
        top = Inches(0.3)
        width = Inches(11.33)  # Increased width for widescreen
        height = Inches(0.8)

        # Main title and subheader
        if slide_type == "visitors_trips":
            title_text = "Canadian Visitation"
            subheader_text = f"Canadian Visitors & Trips Travelling to {dmo_name} by month"
        else:
            title_text = "Canadian Overnight Stays"
            subheader_text = f"Canadian Overnight Stays in {dmo_name} by month"

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True

        subheader_box = slide.shapes.add_textbox(left, top + Inches(0.7), width, Inches(0.6))
        subheader_frame = subheader_box.text_frame
        subheader_frame.text = subheader_text
        subheader_frame.paragraphs[0].font.size = Pt(22)
        subheader_frame.paragraphs[0].font.bold = False

        # Position charts in center area with better spacing for widescreen
        img_left = Inches(2.0)  # Centered horizontally
        img_top = top + Inches(1.5)
        max_img_width = Inches(9.33)  # Maximum width available for charts
        vertical_spacing = Inches(0.1)  # Space between charts

        for idx, (period_type, chart_path) in enumerate(charts):
            # Get the original image dimensions to maintain aspect ratio
            try:
                with Image.open(chart_path) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height

                    # Calculate height based on desired width and original aspect ratio
                    desired_width = max_img_width
                    desired_height = desired_width / aspect_ratio

                    # If too tall, adjust based on height constraint
                    max_height = Inches(2.7)
                    if desired_height > max_height:
                        desired_height = max_height
                        desired_width = desired_height * aspect_ratio

                    # Center the image horizontally
                    adjusted_left = img_left + (max_img_width - desired_width) / 2

                    # Position with vertical spacing
                    chart_top = img_top + idx * (max_height + vertical_spacing)

                    # Add image while preserving aspect ratio
                    slide.shapes.add_picture(chart_path, adjusted_left, chart_top,
                                             width=desired_width, height=desired_height)
            except Exception as e:
                print(f"[ERROR] Failed to process image {chart_path}: {e}")
                # Fallback to previous method if image processing fails
                chart_top = img_top + idx * (Inches(2.7) + vertical_spacing)
                slide.shapes.add_picture(chart_path, img_left, chart_top,
                                         width=max_img_width, height=Inches(2.7))

    print("[DEBUG] Finished adding line chart slides, starting grouped barplot generation")
    # Do NOT generate grouped barplots here; they are already generated in the GUI.
    print("[DEBUG] Finished grouped barplot generation, adding barplot slides")

    # Generate barplots for each DMO using the last DMO name from the slides list
    dmo_names = set()
    for slide_info in slides:
        dmo_name = slide_info["dmo_name"]
        # Use DMO column value as the name if available
        if "DMO" in slide_info:
            dmo_name = slide_info["DMO"]
        dmo_names.add(dmo_name)

    # After all line chart slides, add grouped barplot slides
    # Look for barplot files that might have been generated with quarterly data
    for dmo_name in dmo_names:
        # Convert DMO name to a format that might be used in filenames
        safe_name = dmo_name.replace(' ', '_').replace(',', '').replace('(', '').replace(')', '')

        # Check for both potential naming patterns
        barplot_path = os.path.join('output', f"{safe_name}_top5_origins.png")
        alt_barplot_path = None

        # If the barplot file doesn't exist with the standard name, look for alternative formats
        if not os.path.exists(barplot_path):
            # Look for any file that might contain the DMO name and quarterly data
            potential_barplots = glob.glob(os.path.join('output', f"{safe_name}*_barplot*.png"))
            if not potential_barplots:
                potential_barplots = glob.glob(os.path.join('output', f"{safe_name}*origins*.png"))
            if potential_barplots:
                barplot_path = potential_barplots[0]
            else:
                # Skip this DMO if no barplot is found
                print(f"[DEBUG] No barplot found for {dmo_name}")
                continue

        print(f"[DEBUG] Adding barplot slide for {barplot_path}")

        # Check if there's an out-of-province barplot for this DMO
        outprov_barplot_path = os.path.join('output', f"{safe_name}_top5_outofprov_origins.png")
        # Look for alternate out-of-province barplot naming
        if not os.path.exists(outprov_barplot_path):
            potential_outprov = glob.glob(os.path.join('output', f"{safe_name}*outprov*.png"))
            potential_outprov += glob.glob(os.path.join('output', f"{safe_name}*out-of-prov*.png"))
            if potential_outprov:
                outprov_barplot_path = potential_outprov[0]

        has_outprov_barplot = os.path.exists(outprov_barplot_path)

        slide = prs.slides.add_slide(content_slide_layout)
        # Remove all default placeholders
        remove_placeholders(slide)
        # Add title
        left = Inches(1.0)  # Increased from 0.5
        top = Inches(0.3)
        width = Inches(11.33)  # Increased width for widescreen
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = f"Top 5 Origins Visiting {dmo_name}"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Size and position adjustments for widescreen format with aspect ratio preservation
        if has_outprov_barplot:
            # If we have two barplot charts, position them side by side with pie and nights charts below
            label_top = top + Inches(0.8)
            label_height = Inches(0.4)

            # Add labels for each chart
            # Label for left chart
            label_left = Inches(1.0)
            label_box = slide.shapes.add_textbox(label_left, label_top, MAX_CHART_WIDTH, label_height)
            label_frame = label_box.text_frame
            label_frame.text = "All Origins"
            label_frame.paragraphs[0].font.size = LABEL_FONT_SIZE
            label_frame.paragraphs[0].font.bold = True
            label_frame.paragraphs[0].alignment = 1  # Center alignment

            # Label for right chart
            label_left = Inches(6.83)
            label_box = slide.shapes.add_textbox(label_left, label_top, MAX_CHART_WIDTH, label_height)
            label_frame = label_box.text_frame
            label_frame.text = "Out-of-Province Origins"
            label_frame.paragraphs[0].font.size = LABEL_FONT_SIZE
            label_frame.paragraphs[0].font.bold = True
            label_frame.paragraphs[0].alignment = 1  # Center alignment

            # Left chart with preserved aspect ratio
            add_image_with_aspect_ratio(slide, barplot_path, Inches(1.0), top + Inches(1.2),
                                       MAX_CHART_WIDTH, Inches(2.5))

            # Right chart with preserved aspect ratio
            add_image_with_aspect_ratio(slide, outprov_barplot_path, Inches(6.83), top + Inches(1.2),
                                       MAX_CHART_WIDTH, Inches(2.5))

            # Check for pie chart and nights per visitor chart paths
            pie_chart_path = os.path.join('output', f"{safe_name}_province_pie.png")
            nights_per_visitor_path = os.path.join('output', f"{safe_name}_nights_per_visitor_quarterly.png")

            # If pie chart doesn't exist with standard name, look for alternatives
            if not os.path.exists(pie_chart_path):
                potential_pie_charts = glob.glob(os.path.join('output', f"{safe_name}*_pie*.png"))
                if potential_pie_charts:
                    pie_chart_path = potential_pie_charts[0]
                    print(f"[DEBUG] Found alternative pie chart path: {pie_chart_path}")
                else:
                    pie_chart_path = None
                    print(f"[DEBUG] No province pie chart found for {dmo_name}")

            # Check for nights per visitor chart
            if not os.path.exists(nights_per_visitor_path):
                potential_nights_charts = glob.glob(os.path.join('output', f"{safe_name}*nights_per_visitor*.png"))
                if potential_nights_charts:
                    nights_per_visitor_path = potential_nights_charts[0]
                    print(f"[DEBUG] Found nights per visitor chart: {nights_per_visitor_path}")
                else:
                    nights_per_visitor_path = None
                    print(f"[DEBUG] No nights per visitor chart found for {dmo_name}")

            # Add charts at the bottom in a row
            bottom_charts_top = top + Inches(4.0)

            # Add nights per visitor chart on the left
            if nights_per_visitor_path and os.path.exists(nights_per_visitor_path):
                # Add label for nights per visitor chart
                nights_label_left = Inches(1.0)
                nights_label_top = bottom_charts_top - Inches(0.3)
                nights_label_box = slide.shapes.add_textbox(nights_label_left, nights_label_top,
                                                          Inches(5.0), Inches(0.3))
                nights_label_frame = nights_label_box.text_frame
                nights_label_frame.text = "Nights per Visitor by Quarter"
                nights_label_frame.paragraphs[0].font.size = LABEL_FONT_SIZE
                nights_label_frame.paragraphs[0].font.bold = True
                nights_label_frame.paragraphs[0].alignment = 1  # Center alignment

                add_image_with_aspect_ratio(slide, nights_per_visitor_path, Inches(1.0), bottom_charts_top,
                                           Inches(5.0), Inches(2.5))

            # Add province pie chart on the right
            if pie_chart_path and os.path.exists(pie_chart_path):
                # Add label for pie chart
                pie_label_left = Inches(7.0)
                pie_label_top = bottom_charts_top - Inches(0.3)
                pie_label_box = slide.shapes.add_textbox(pie_label_left, pie_label_top,
                                                       Inches(4.0), Inches(0.3))
                pie_label_frame = pie_label_box.text_frame
                pie_label_frame.text = "Visitors by Province of Origin"
                pie_label_frame.paragraphs[0].font.size = LABEL_FONT_SIZE
                pie_label_frame.paragraphs[0].font.bold = True
                pie_label_frame.paragraphs[0].alignment = 1  # Center alignment

                add_image_with_aspect_ratio(slide, pie_chart_path, Inches(7.0), bottom_charts_top,
                                           Inches(4.0), Inches(2.5))
        else:
            # If we only have one barplot chart, arrange all three charts nicely
            # Add label for barplot
            barplot_label_width = Inches(6.5)
            barplot_label_height = Inches(0.4)
            barplot_label_left = Inches(1.0)
            barplot_label_top = top + Inches(0.8)

            barplot_label_box = slide.shapes.add_textbox(barplot_label_left, barplot_label_top,
                                                       barplot_label_width, barplot_label_height)
            barplot_label_frame = barplot_label_box.text_frame
            barplot_label_frame.text = "All Origins"
            barplot_label_frame.paragraphs[0].font.size = LABEL_FONT_SIZE
            barplot_label_frame.paragraphs[0].font.bold = True
            barplot_label_frame.paragraphs[0].alignment = 1  # Center alignment

            # Add barplot at the top
            add_image_with_aspect_ratio(slide, barplot_path, Inches(0.8), top + Inches(1.2),
                                       Inches(10.0), Inches(3.0))

            # Check for pie chart and nights per visitor chart paths
            pie_chart_path = os.path.join('output', f"{safe_name}_province_pie.png")
            nights_per_visitor_path = os.path.join('output', f"{safe_name}_nights_per_visitor_quarterly.png")

            # If pie chart doesn't exist with standard name, look for alternatives
            if not os.path.exists(pie_chart_path):
                potential_pie_charts = glob.glob(os.path.join('output', f"{safe_name}*_pie*.png"))
                if potential_pie_charts:
                    pie_chart_path = potential_pie_charts[0]
                    print(f"[DEBUG] Found alternative pie chart path: {pie_chart_path}")
                else:
                    pie_chart_path = None
                    print(f"[DEBUG] No province pie chart found for {dmo_name}")

            # Check for nights per visitor chart
            if not os.path.exists(nights_per_visitor_path):
                potential_nights_charts = glob.glob(os.path.join('output', f"{safe_name}*nights_per_visitor*.png"))
                if potential_nights_charts:
                    nights_per_visitor_path = potential_nights_charts[0]
                    print(f"[DEBUG] Found nights per visitor chart: {nights_per_visitor_path}")
                else:
                    nights_per_visitor_path = None
                    print(f"[DEBUG] No nights per visitor chart found for {dmo_name}")

            # Add charts at the bottom in a row
            bottom_charts_top = top + Inches(4.5)

            # Add nights per visitor chart on the left
            if nights_per_visitor_path and os.path.exists(nights_per_visitor_path):
                # Add label for nights per visitor chart
                nights_label_left = Inches(1.0)
                nights_label_top = bottom_charts_top - Inches(0.3)
                nights_label_box = slide.shapes.add_textbox(nights_label_left, nights_label_top,
                                                          Inches(5.0), Inches(0.3))
                nights_label_frame = nights_label_box.text_frame
                nights_label_frame.text = "Nights per Visitor by Quarter"
                nights_label_frame.paragraphs[0].font.size = LABEL_FONT_SIZE
                nights_label_frame.paragraphs[0].font.bold = True
                nights_label_frame.paragraphs[0].alignment = 1  # Center alignment

                add_image_with_aspect_ratio(slide, nights_per_visitor_path, Inches(1.0), bottom_charts_top,
                                           Inches(5.0), Inches(2.2))

            # Add pie chart on the right
            if pie_chart_path and os.path.exists(pie_chart_path):
                # Add label for pie chart
                pie_label_left = Inches(7.0)
                pie_label_top = bottom_charts_top - Inches(0.3)
                pie_label_box = slide.shapes.add_textbox(pie_label_left, pie_label_top,
                                                       Inches(4.0), Inches(0.3))
                pie_label_frame = pie_label_box.text_frame
                pie_label_frame.text = "Visitors by Province of Origin"
                pie_label_frame.paragraphs[0].font.size = LABEL_FONT_SIZE
                pie_label_frame.paragraphs[0].font.bold = True
                pie_label_frame.paragraphs[0].alignment = 1  # Center alignment

                add_image_with_aspect_ratio(slide, pie_chart_path, Inches(7.0), bottom_charts_top,
                                           Inches(4.0), Inches(2.2))

    # Add PRIZM & Traveller Segmentation Program slides for each DMO destination
    year_from_slides = slides[0]["year"] if slides else None
    comp_year_from_slides = slides[0]["comp_year"] if slides and slides[0]["comp_year"] else None

    # Create individual PRIZM slides for each DMO destination
    for dmo_name in dmo_names:
        # Convert DMO name to a format that might be used in filenames
        safe_name = dmo_name.replace(' ', '_').replace(',', '').replace('(', '').replace(')', '')

        # Look for circular barplot files for this specific DMO
        main_year_circular_path = os.path.join('output', f"{safe_name}_prizm_circular_barplot_{year_from_slides}.png")
        comp_year_circular_path = None
        if comp_year_from_slides:
            comp_year_circular_path = os.path.join('output', f"{safe_name}_prizm_circular_barplot_{comp_year_from_slides}.png")

        # Look for ordinary PRIZM barplot files for this specific DMO
        main_year_prizm_path = os.path.join('output', f"{safe_name}_prizm_barplot_{year_from_slides}.png")
        comp_year_prizm_path = None
        if comp_year_from_slides:
            comp_year_prizm_path = os.path.join('output', f"{safe_name}_prizm_barplot_{comp_year_from_slides}.png")

        # Only add the slide if we have at least the main year chart for this DMO
        if os.path.exists(main_year_circular_path):
            print(f"[DEBUG] Adding PRIZM & Traveller Segmentation slide for {dmo_name}")

            slide = prs.slides.add_slide(content_slide_layout)
            remove_placeholders(slide)

            # Add title
            left = Inches(1.0)
            top = Inches(0.3)
            width = Inches(11.33)
            height = Inches(0.8)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame

            if comp_year_from_slides and comp_year_circular_path and os.path.exists(comp_year_circular_path):
                title_frame.text = f"{dmo_name} - {comp_year_from_slides} & {year_from_slides} Visitors by PRIZM & Traveller Segmentation Program"
            else:
                title_frame.text = f"{dmo_name} - {year_from_slides} Visitors by PRIZM & Traveller Segmentation Program"

            title_frame.paragraphs[0].font.size = Pt(24)  # Slightly smaller to fit longer title
            title_frame.paragraphs[0].font.bold = True

            # Position charts in a 2x2 grid layout
            chart_top_row = top + Inches(1.0)
            chart_bottom_row = top + Inches(4.2)
            chart_left_col = Inches(1.0)
            chart_right_col = Inches(7.0)
            chart_width = Inches(5.5)
            chart_height = Inches(2.8)

            if comp_year_from_slides and comp_year_circular_path and os.path.exists(comp_year_circular_path):
                # Four charts in 2x2 grid
                # Top row: Circular barplots (EQ segments)
                # Bottom row: Ordinary barplots (PRIZM segments)

                # Top left: Comparison year circular barplot
                add_image_with_aspect_ratio(slide, comp_year_circular_path, chart_left_col, chart_top_row,
                                           chart_width, chart_height)
                # Add label
                left_label_box = slide.shapes.add_textbox(chart_left_col, chart_top_row - Inches(0.4), chart_width, Inches(0.3))
                left_label_frame = left_label_box.text_frame
                left_label_frame.text = f"{comp_year_from_slides} - Traveller Segmentation"
                left_label_frame.paragraphs[0].font.size = Pt(14)
                left_label_frame.paragraphs[0].font.bold = True
                left_label_frame.paragraphs[0].alignment = 1  # Center alignment

                # Top right: Main year circular barplot
                add_image_with_aspect_ratio(slide, main_year_circular_path, chart_right_col, chart_top_row,
                                           chart_width, chart_height)
                # Add label
                right_label_box = slide.shapes.add_textbox(chart_right_col, chart_top_row - Inches(0.4), chart_width, Inches(0.3))
                right_label_frame = right_label_box.text_frame
                right_label_frame.text = f"{year_from_slides} - Traveller Segmentation"
                right_label_frame.paragraphs[0].font.size = Pt(14)
                right_label_frame.paragraphs[0].font.bold = True
                right_label_frame.paragraphs[0].alignment = 1  # Center alignment

                # Bottom left: Comparison year ordinary PRIZM barplot
                if comp_year_prizm_path and os.path.exists(comp_year_prizm_path):
                    add_image_with_aspect_ratio(slide, comp_year_prizm_path, chart_left_col, chart_bottom_row,
                                               chart_width, chart_height)
                    # Add label
                    bottom_left_label_box = slide.shapes.add_textbox(chart_left_col, chart_bottom_row - Inches(0.4), chart_width, Inches(0.3))
                    bottom_left_label_frame = bottom_left_label_box.text_frame
                    bottom_left_label_frame.text = f"{comp_year_from_slides} - PRIZM Segments"
                    bottom_left_label_frame.paragraphs[0].font.size = Pt(14)
                    bottom_left_label_frame.paragraphs[0].font.bold = True
                    bottom_left_label_frame.paragraphs[0].alignment = 1  # Center alignment

                # Bottom right: Main year ordinary PRIZM barplot
                if main_year_prizm_path and os.path.exists(main_year_prizm_path):
                    add_image_with_aspect_ratio(slide, main_year_prizm_path, chart_right_col, chart_bottom_row,
                                               chart_width, chart_height)
                    # Add label
                    bottom_right_label_box = slide.shapes.add_textbox(chart_right_col, chart_bottom_row - Inches(0.4), chart_width, Inches(0.3))
                    bottom_right_label_frame = bottom_right_label_box.text_frame
                    bottom_right_label_frame.text = f"{year_from_slides} - PRIZM Segments"
                    bottom_right_label_frame.paragraphs[0].font.size = Pt(14)
                    bottom_right_label_frame.paragraphs[0].font.bold = True
                    bottom_right_label_frame.paragraphs[0].alignment = 1  # Center alignment
            else:
                # Two charts side by side (only main year available)
                # Left: Circular barplot (EQ segments)
                # Right: Ordinary barplot (PRIZM segments)

                chart_width_single = Inches(5.5)
                chart_height_single = Inches(5.0)
                chart_top_single = top + Inches(1.5)

                # Left: Main year circular barplot
                add_image_with_aspect_ratio(slide, main_year_circular_path, chart_left_col, chart_top_single,
                                           chart_width_single, chart_height_single)
                # Add label
                left_single_label_box = slide.shapes.add_textbox(chart_left_col, chart_top_single - Inches(0.4), chart_width_single, Inches(0.3))
                left_single_label_frame = left_single_label_box.text_frame
                left_single_label_frame.text = f"{year_from_slides} - Traveller Segmentation"
                left_single_label_frame.paragraphs[0].font.size = Pt(14)
                left_single_label_frame.paragraphs[0].font.bold = True
                left_single_label_frame.paragraphs[0].alignment = 1  # Center alignment

                # Right: Main year ordinary PRIZM barplot
                if main_year_prizm_path and os.path.exists(main_year_prizm_path):
                    add_image_with_aspect_ratio(slide, main_year_prizm_path, chart_right_col, chart_top_single,
                                               chart_width_single, chart_height_single)
                    # Add label
                    right_single_label_box = slide.shapes.add_textbox(chart_right_col, chart_top_single - Inches(0.4), chart_width_single, Inches(0.3))
                    right_single_label_frame = right_single_label_box.text_frame
                    right_single_label_frame.text = f"{year_from_slides} - PRIZM Segments"
                    right_single_label_frame.paragraphs[0].font.size = Pt(14)
                    right_single_label_frame.paragraphs[0].font.bold = True
                    right_single_label_frame.paragraphs[0].alignment = 1  # Center alignment
        else:
            print(f"[DEBUG] No PRIZM circular barplot found for {dmo_name} in year {year_from_slides}")

    # Add Quarterly PRIZM slides for each DMO destination (if comparison year exists)
    if comp_year_from_slides:
        for dmo_name in dmo_names:
            # Convert DMO name to a format that might be used in filenames
            safe_name = dmo_name.replace(' ', '_').replace(',', '').replace('(', '').replace(')', '')

            # Look for quarterly PRIZM barplot files for this specific DMO
            quarterly_prizm_path = os.path.join('output', f"{safe_name}_quarterly_prizm_barplots_{comp_year_from_slides}_{year_from_slides}.png")

            # Only add the slide if we have the quarterly chart for this DMO
            if os.path.exists(quarterly_prizm_path):
                print(f"[DEBUG] Adding Quarterly PRIZM slide for {dmo_name}")

                slide = prs.slides.add_slide(content_slide_layout)
                remove_placeholders(slide)

                # Add title
                left = Inches(1.0)
                top = Inches(0.3)
                width = Inches(11.33)
                height = Inches(0.8)
                title_box = slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_frame.text = f"{dmo_name} - Quarterly {comp_year_from_slides} & {year_from_slides} Visitors by PRIZM"
                title_frame.paragraphs[0].font.size = Pt(28)
                title_frame.paragraphs[0].font.bold = True

                # Add the quarterly chart (single large image)
                chart_left = Inches(0.5)
                chart_top = top + Inches(1.2)
                chart_width = Inches(12.33)
                chart_height = Inches(5.5)

                add_image_with_aspect_ratio(slide, quarterly_prizm_path, chart_left, chart_top,
                                           chart_width, chart_height)

                # Add description text below the chart
                description_top = chart_top + chart_height + Inches(0.2)
                description_box = slide.shapes.add_textbox(left, description_top, width, Inches(0.6))
                description_frame = description_box.text_frame
                description_frame.text = f"Top row: {comp_year_from_slides} quarters (Q1-Q4) • Bottom row: {year_from_slides} quarters (Q1-Q4) • Each chart shows top 5 PRIZM segments for that quarter"
                description_frame.paragraphs[0].font.size = Pt(14)
                description_frame.paragraphs[0].font.bold = False
                description_frame.paragraphs[0].alignment = 1  # Center alignment
            else:
                print(f"[DEBUG] No quarterly PRIZM barplot found for {dmo_name}")

    print("[DEBUG] Attempting to save PowerPoint file")
    try:
        prs.save(ppt_path)
        print(f"[DEBUG] PowerPoint saved to {ppt_path}")
    except Exception as e:
        print(f"Error saving PowerPoint file to {ppt_path}: {e}")
        raise
    print("[DEBUG] create_presentation complete")
